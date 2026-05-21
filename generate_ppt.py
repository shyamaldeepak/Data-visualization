from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Slide layouts
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
title_only_layout = prs.slide_layouts[5]

# Slide 1: Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Video Game Sales Analysis"
subtitle.text = "Group Data Visualization Project\nMember 1 | Member 2 | Member 3 | Member 4 | Member 5"

# Slide 2: Context & Data (Member 1)
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
title.text = "Project Context & Data (Member 1)"
tf = slide.placeholders[1].text_frame
tf.text = "Dataset: Video Game Sales (Kaggle / VGChartz, Oct 2016)"
p = tf.add_paragraph()
p.text = "Records: 16,598 games, 11 attributes"
p.level = 1
p = tf.add_paragraph()
p.text = "Main Question: What are the primary factors that drive the commercial success of a video game globally?"

# Slide 3: Graph 1 - Genres (Member 1)
slide = prs.slides.add_slide(title_only_layout)
title = slide.shapes.title
title.text = "1. Top Selling Genres"
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
tf = txBox.text_frame
tf.text = "=> Action and Sports games dominate global sales due to broad demographic appeal.\n[INSERT GENRE BAR CHART HERE]"

# Slide 4: Graph 2 - Platforms (Member 2)
slide = prs.slides.add_slide(title_only_layout)
title = slide.shapes.title
title.text = "2. Top Performing Platforms (Member 2)"
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
tf = txBox.text_frame
tf.text = "=> PS2, Xbox 360, and Wii lead historical sales, reflecting the 2000s console boom.\n[INSERT PLATFORM BAR CHART HERE]"

# Slide 5: Graph 3 - Temporal Trends (Member 3)
slide = prs.slides.add_slide(title_only_layout)
title = slide.shapes.title
title.text = "3. Sales Over Time (Member 3)"
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
tf = txBox.text_frame
tf.text = "=> Sales peaked heavily around 2008-2009. The drop after 2010 is largely because digital sales are missing from this dataset.\n[INSERT YEARS LINE CHART HERE]"

# Slide 6: Graph 4 - Regional Market Share (Member 4)
slide = prs.slides.add_slide(title_only_layout)
title = slide.shapes.title
title.text = "4. Regional Contributions (Member 4)"
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
tf = txBox.text_frame
tf.text = "=> North America is the dominant market (~50% of global sales).\n[INSERT REGIONAL PIE/BAR CHART HERE]"

# Slide 7: Graph 5 - Correlation (Member 5)
slide = prs.slides.add_slide(title_only_layout)
title = slide.shapes.title
title.text = "5. Rank vs. Sales Relationship (Member 5)"
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
tf = txBox.text_frame
tf.text = "=> Correlation: -0.42. The absolute top 10 ranked games sell exponentially more than the rest.\n[INSERT SCATTER PLOT HERE]"

# Slide 8: Conclusion & Critical Thinking (Member 5)
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
title.text = "Conclusions & Critical Thinking"
tf = slide.placeholders[1].text_frame
tf.text = "Key Drivers: Action/Sports genres, major consoles (PS2/Wii), NA market."
p = tf.add_paragraph()
p.text = "Limitations: Data ends in 2016 and lacks digital sales."
p.level = 0
p = tf.add_paragraph()
p.text = "Next Steps: Incorporate modern digital APIs and apply machine learning models."
p.level = 0

# Slide 9: Q&A (Member 5)
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
title.text = "Questions & Answers"
tf = slide.placeholders[1].text_frame
tf.text = "Question for the Audience: Do you think Action games still dominate today, or have mobile games completely changed the landscape?"

prs.save('Video_Game_Sales_Presentation.pptx')
